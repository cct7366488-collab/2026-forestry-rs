"""
產生：03-土肉桂專區說明會簡報.pptx
為 6/6 土肉桂專區系統說明會準備（對應 prod v2.11.69）。

16:9 widescreen 模式、無講稿備註（user 自己股語）、緊扣 6 段演示腳本。

跑法：python _build_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image as _PILImage

HERE = Path(__file__).parent
OUT = HERE / "03-土肉桂專區說明會簡報.pptx"
IMG = HERE / "images"

# 色彩主題（與 ForestMRV 一致：森林綠）
FOREST_700 = RGBColor(0x15, 0x80, 0x3D)
FOREST_800 = RGBColor(0x16, 0x65, 0x34)
FOREST_900 = RGBColor(0x14, 0x53, 0x2D)
STONE_700 = RGBColor(0x44, 0x40, 0x3B)
STONE_500 = RGBColor(0x78, 0x71, 0x69)
STONE_300 = RGBColor(0xD6, 0xD3, 0xD1)
STONE_100 = RGBColor(0xF5, 0xF5, 0xF4)
AMBER_500 = RGBColor(0xF5, 0x9E, 0x0B)
BLUE_700 = RGBColor(0x1D, 0x4E, 0xD8)
RED_600 = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CJK_FONT = "Noto Sans TC"
ASCII_FONT = "Calibri"

prs = Presentation()
# 16:9 widescreen (13.333" x 7.5")
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_text(run, size=18, bold=False, italic=False, color=None, font_name=CJK_FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    run.font.name = font_name
    # 設定東亞字型
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    eaFont = rPr.find(qn('a:ea'))
    if eaFont is None:
        eaFont = etree.SubElement(rPr, qn('a:ea'))
    eaFont.set('typeface', CJK_FONT)


IMG_BG_COVER = IMG / "bg-cover.png"
IMG_BG_SECTION = IMG / "bg-section.png"


def _set_solid_alpha(shape, opacity_pct):
    """在 solidFill 的 srgbClr 注入 <a:alpha>，opacity_pct=不透明度(0-100)"""
    srgb = shape._element.spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    a = etree.SubElement(srgb, qn('a:alpha'))
    a.set('val', str(int(opacity_pct * 1000)))


def add_full_bleed_bg(slide, img_path, overlay_color, overlay_opacity, fallback_color):
    """滿版底圖 + 半透明遮罩；圖檔不存在則退回純色底。回傳是否用了底圖。"""
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        ov = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                    prs.slide_width, prs.slide_height)
        ov.fill.solid()
        ov.fill.fore_color.rgb = overlay_color
        ov.line.fill.background()
        ov.shadow.inherit = False
        _set_solid_alpha(ov, overlay_opacity)
        return True
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = fallback_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return False


def add_node(slide, left, top, w, h, lines, fill=FOREST_700,
             line_color=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """流程示意圖節點：圓角矩形 + 置中多行文字。lines = [(text,size,bold,color),...]"""
    shp = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, sz, bold, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        set_text(r, size=sz, bold=bold, color=col)
    return shp


def add_arrow(slide, left, top, w, h, color=AMBER_500, direction="right"):
    sh = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    a = slide.shapes.add_shape(sh, Inches(left), Inches(top), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def fit_in_box(img_path, max_w_in, max_h_in):
    """等比例縮放使圖片同時不超過 max_w × max_h（英吋），回傳 (w_in, h_in)。"""
    with _PILImage.open(img_path) as im:
        w_px, h_px = im.size
    ar = h_px / w_px
    w = max_w_in
    h = w * ar
    if h > max_h_in:
        h = max_h_in
        w = h / ar
    return w, h


def new_diagram_slide(title):
    """空白頁 + 標題條 + 頁尾，供繪向量示意圖。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, title)
    add_footer(slide)
    return slide


def add_notes_box(slide, left, top, w, h, items, size=20):
    """示意圖下方／旁的補充說明（# 開頭為小標）。"""
    from pptx.enum.text import MSO_AUTO_SIZE
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        if it.startswith("# "):
            r = p.add_run()
            r.text = it[2:]
            set_text(r, size=size + 5, bold=True, color=FOREST_700)
        else:
            r = p.add_run()
            r.text = "• " + it
            set_text(r, size=size, color=STONE_700)


def add_title_bar(slide, title, subtitle=None):
    """加頂部標題條（綠底白字）"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOREST_800
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_text(run, size=34, bold=True, color=WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        run2 = p2.add_run()
        run2.text = subtitle
        set_text(run2, size=18, color=STONE_100)


def add_footer(slide, text="ForestMRV v2.11.69 │ 土肉桂專區葉片收穫監測會計系統"):
    """加底部頁尾"""
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(13), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_text(run, size=11, color=STONE_500)


def add_section_divider(prs, section_no, section_title, time_range=None):
    """章節分隔頁（綠底大字置中）"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_full_bleed_bg(slide, IMG_BG_SECTION, FOREST_900, 62, FOREST_800)

    # 段號
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.9), Inches(10.3), Inches(1.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"第 {section_no} 段"
    set_text(run, size=48, color=AMBER_500, bold=True)

    # 段標題
    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(3.3), Inches(11.3), Inches(2.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = section_title
    set_text(run2, size=58, color=WHITE, bold=True)

    if time_range:
        tb3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.9))
        tf3 = tb3.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = time_range
        set_text(run3, size=27, color=STONE_300)


def add_content_slide(prs, title, bullets=None, image=None, image_caption=None,
                      image_width=5.5, image_left=7.0, image_top=1.4):
    """內容頁：標題 + 條列 + 可選右側圖"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, title)

    # 條列區（左半邊或全寬，依是否有圖）
    text_width = Inches(6.3) if image else Inches(12.5)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), text_width, Inches(5.6))
    tf = tb.text_frame
    tf.margin_left = Inches(0.1)
    tf.word_wrap = True
    # 字級放大後的安全網：內容溢出時自動縮字（PowerPoint normAutofit）
    from pptx.enum.text import MSO_AUTO_SIZE
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    if bullets:
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(5)
            p.space_after = Pt(5)
            # 處理 (heading, body) tuple 或單純 string
            if isinstance(b, tuple):
                heading, body = b
                run = p.add_run()
                run.text = heading
                set_text(run, size=24, bold=True, color=FOREST_700)
                if body:
                    run2 = p.add_run()
                    run2.text = " " + body
                    set_text(run2, size=20, color=STONE_700)
            elif b.startswith("# "):
                run = p.add_run()
                run.text = b[2:]
                set_text(run, size=27, bold=True, color=FOREST_700)
            elif b.startswith("> "):
                p.level = 1
                run = p.add_run()
                run.text = " " + b[2:]
                set_text(run, size=19, italic=True, color=STONE_500)
            else:
                run = p.add_run()
                run.text = "• " + b
                set_text(run, size=20, color=STONE_700)

    if image:
        img_path = IMG / image
        if img_path.exists():
            # 面板可用範圍：寬 image_width，高至頁尾上方（留圖說空間）
            panel_w = image_width
            max_h = 6.75 - image_top - (0.35 if image_caption else 0.0)
            fw, fh = fit_in_box(img_path, panel_w, max_h)
            # 在面板水平置中
            cx_left = image_left + (panel_w - fw) / 2
            slide.shapes.add_picture(str(img_path), Inches(cx_left), Inches(image_top),
                                     width=Inches(fw), height=Inches(fh))
            if image_caption:
                cap_tb = slide.shapes.add_textbox(
                    Inches(image_left), Inches(image_top + fh + 0.08),
                    Inches(image_width), Inches(0.3))
                cap_tf = cap_tb.text_frame
                cap_tf.margin_left = 0
                cp = cap_tf.paragraphs[0]
                cp.alignment = PP_ALIGN.CENTER
                crun = cp.add_run()
                crun.text = image_caption
                set_text(crun, size=13, italic=True, color=STONE_500)
        else:
            # 圖檔缺失提示
            tb_err = slide.shapes.add_textbox(Inches(image_left), Inches(image_top),
                                              Inches(image_width), Inches(1))
            tf_err = tb_err.text_frame
            p = tf_err.paragraphs[0]
            run = p.add_run()
            run.text = f"［圖檔遺失：{image}］"
            set_text(run, size=14, italic=True, color=RED_600)

    add_footer(slide)


def add_big_image_slide(prs, title, image, image_caption=None, image_width=10.0):
    """大圖頁：標題 + 置中大圖（無條列）"""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title_bar(slide, title)
    img_path = IMG / image
    if img_path.exists():
        # 可用範圍：標題條下方 1.25 至圖說上方 6.45，寬 image_width
        region_top, region_bot = 1.25, 6.45
        fw, fh = fit_in_box(img_path, image_width, region_bot - region_top)
        left = (prs.slide_width - Inches(fw)) // 2
        top = Inches(region_top + (region_bot - region_top - fh) / 2)
        slide.shapes.add_picture(str(img_path), left, top,
                                 width=Inches(fw), height=Inches(fh))
        if image_caption:
            cap_tb = slide.shapes.add_textbox(Inches(1), Inches(6.55),
                                               Inches(11.3), Inches(0.4))
            cap_tf = cap_tb.text_frame
            cap_tf.margin_left = 0
            cp = cap_tf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            crun = cp.add_run()
            crun.text = image_caption
            set_text(crun, size=14, italic=True, color=STONE_500)
    else:
        tb_err = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.3), Inches(0.5))
        tf_err = tb_err.text_frame
        p = tf_err.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"［圖檔遺失：{image}］"
        set_text(run, size=18, italic=True, color=RED_600)
    add_footer(slide)


# =====================================================
# Slide 1: 封面
# =====================================================
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)
add_full_bleed_bg(slide, IMG_BG_COVER, FOREST_900, 58, FOREST_900)

# 主標題
tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "土肉桂專區"
set_text(run, size=72, bold=True, color=AMBER_500)

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "葉片收穫監測會計系統"
set_text(run2, size=54, bold=True, color=WHITE)

# 副標題
tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.0))
tf_sub = tb_sub.text_frame
p_sub = tf_sub.paragraphs[0]
p_sub.alignment = PP_ALIGN.CENTER
run_sub = p_sub.add_run()
run_sub.text = "系統說明會"
set_text(run_sub, size=36, color=STONE_100)

# 底部資訊
tb_info = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8))
tf_info = tb_info.text_frame
p_info = tf_info.paragraphs[0]
p_info.alignment = PP_ALIGN.CENTER
run_info = p_info.add_run()
run_info.text = "2026 年 6 月 6 日 │ ForestMRV v2.11.69 │ 陳朝圳"
set_text(run_info, size=20, color=STONE_300)

# =====================================================
# Slide 2: 議程
# =====================================================
add_content_slide(prs, "議程（30 分鐘）",
    bullets=[
        ("第 1 段（0-5 分鐘）", "系統概觀與角色介紹"),
        ("第 2 段（5-10 分鐘）", "系統管理員開案與邊界準備"),
        ("第 3 段（10-15 分鐘）", "林農修枝申請（主秀）"),
        ("第 4 段（15-20 分鐘）", "分署承辦審核與核准"),
        ("第 5 段（20-25 分鐘）", "林農回報葉片採收量與結案"),
        ("第 6 段（25-30 分鐘）", "合作社彙整與 Q&A"),
    ])

# =====================================================
# Slide 3: 系統定位
# =====================================================
add_content_slide(prs, "系統定位",
    bullets=[
        "# ForestMRV 智慧森林監測平臺之子模組",
        "依循 IPCC LULUCF 與 MRV（測量、報告、查證）原則",
        "土肉桂專區採「行政管理導向」設計、非樣區調查專案",
        "# 承接四階段行政流程",
        "1. 林農申請修枝採取副產物（葉片）",
        "2. 林業保育署臺中分署核准並核發法定許可文號",
        "3. 林農回報實際葉片採收量",
        "4. 結案後資料移交合作社進行共同銷售彙整",
    ])

# =====================================================
# Slide 4: 五個角色
# =====================================================
add_content_slide(prs, "五個系統角色",
    bullets=[
        ("admin（系統管理員）", "計畫團隊 │ 開案、邊界、人員指派、QAQC 監督"),
        ("pi（計畫主持人）", "計畫團隊／合作社主席（移交後）│ 方法學設定、邀請成員"),
        ("surveyor（林農）", "契約承租林農 │ 修枝申請、葉片採收回報、結案"),
        ("harvest_authority（審核人員）", "臺中分署承辦 │ 核准、要求補件、駁回、核發法定許可文號"),
        ("coop（合作社）", "合作社主席 │ 唯讀彙整、共同銷售匯出 Excel"),
        "",
        "> 權限以 Firestore Security Rules 強制執行、介面層 UI gating 為輔助",
    ])

# =====================================================
# Slide 5: 系統入口
# =====================================================
add_content_slide(prs, "系統入口與登入",
    bullets=[
        "# 網址",
        "https://forestry-rs-monitor.web.app",
        "",
        "# 兩種登入方式",
        "Google 登入（建議，免另記密碼）",
        "Email/密碼登入",
        "",
        "# PWA 設計",
        "無需安裝任何軟體",
        "Chrome、Edge、Safari 皆可使用",
        "手機可加到主畫面、體驗近 native app",
    ],
    image="S01-登入頁.png",
    image_caption="ForestMRV 登入頁",
    image_width=4.5, image_left=8.2, image_top=1.4)

# =====================================================
# 第 1 段 分隔頁
# =====================================================
add_section_divider(prs, "1", "系統概觀與角色介紹", "0 - 5 分鐘")

# =====================================================
# Slide 7: 我的專案
# =====================================================
add_big_image_slide(prs, "我的專案首頁",
    image="S02-我的專案.png",
    image_caption="登入後第一個畫面 — admin 可看到全部專案",
    image_width=10.0)

# =====================================================
# Slide 8: 專案首頁與七分頁
# =====================================================
add_content_slide(prs, "專案首頁與七個分頁",
    bullets=[
        "# 純行政專案：分頁已精簡",
        "土肉桂專區不需要樣區、立木、QAQC 等模組",
        "系統自動隱藏「樣區／設計／儀表板」三個分頁",
        "# 七個業務分頁",
        "🗺️ 地圖",
        "🌿 修枝申請",
        "📋 修枝審核",
        "🌾 葉片採收回報及結案",
        "📊 葉片採收彙整",
        "📥 匯出",
        "⚙️ 設定",
    ],
    image="S03-土肉桂專案首頁.png",
    image_caption="土肉桂專案首頁（地圖分頁 + 七個業務分頁）",
    image_width=6.0, image_left=6.8, image_top=1.4)

# =====================================================
# 第 2 段 分隔頁
# =====================================================
add_section_divider(prs, "2", "系統管理員開案與邊界準備", "5 - 10 分鐘")

# =====================================================
# Slide 10: 新建專案
# =====================================================
add_content_slide(prs, "新建土肉桂專案",
    bullets=[
        "# 入口",
        "我的專案 → ＋ 新專案",
        "# 表單重點",
        "專案名稱與編號",
        "計畫類型：「行政管理」或「林產物採取監測」",
        "方法學模組：勾選 admin_harvest",
        ("📐 專案邊界", "點「📥 預設邊界」下拉、選土肉桂專區（橫流溪+烏石坑）"),
        "# 一鍵載入合併邊界",
        "115 個作業單元（橫流溪 45 + 烏石坑 70）",
        "含完整契約 metadata（承租人、契約樹種、面積、契約書號）",
    ],
    image="S15-開案表單.png",
    image_caption="開案表單預設邊界區塊",
    image_width=5.5, image_left=7.3, image_top=1.4)

# =====================================================
# Slide 11: 地圖瀏覽 115 作業單元
# =====================================================
add_big_image_slide(prs, "地圖瀏覽：115 個作業單元",
    image="S17-地圖頁115作業單元.png",
    image_caption="點 label/polygon 跳出 7 欄結構化 metadata popup（作業單元/專區/承租人/契約樹種/面積/工作站/契約書）",
    image_width=10.5)

# =====================================================
# 第 3 段 分隔頁
# =====================================================
add_section_divider(prs, "3", "林農修枝申請（主秀）", "10 - 15 分鐘")

# =====================================================
# Slide 13: 申請清單概觀
# =====================================================
add_content_slide(prs, "修枝申請分頁概觀",
    bullets=[
        "# 業務流程",
        "林農提出申請 → 列印公文寄分署 → 等待核准 → 列印許可單 → 修枝採葉 → 回報採收量 → 結案",
        "# 狀態徽章",
        "草稿 / 已送出．待審 / 已核准 / 修枝作業中 / 已結案 / 已駁回",
        "# 卡片資訊",
        "申請人、林地、面積、株數、修枝期間",
        "✓ 系統識別摘要（picker 帶入結構化資料）",
        "操作按鈕：✏️ 編輯 / 💚 送出申請 / 🗑️ 刪除 / 📄 申請公文稿",
    ],
    image="S04-修枝申請清單.png",
    image_caption="修枝申請清單（含三筆已結案範例）",
    image_width=6.0, image_left=6.8, image_top=1.4)

# =====================================================
# Slide 14: 新增申請表單
# =====================================================
add_content_slide(prs, "新增修枝申請 — 表單設計",
    bullets=[
        "# 表單欄位順序（v2.11.68 設計）",
        "1. 作業位置 picker（三層連動）",
        "2. 作業位置預覽（嵌入迷你地圖）",
        "3. 申請人姓名（picker 自動帶承租人）",
        "4. 聯絡方式（必填）",
        "5. 林班 / 地號（picker 自動帶入）",
        "6. 申請修枝面積（picker 自動帶入）",
        "7. 土肉桂修枝株數（自動 × 1800 株/ha）",
        "8. 修枝方式、起迄日（必填）、備註",
        "# 設計理念",
        "「先選作業位置 → 其他欄位自動帶入 → 補必填」",
    ],
    image="S05-新增申請-空白表單.png",
    image_caption="新增申請表單（picker + 迷你地圖在最上）",
    image_width=4.5, image_left=8.2, image_top=1.4)

# =====================================================
# Slide 15: Picker 三層連動 + 自動帶入
# =====================================================
add_content_slide(prs, "三層 picker 連動 + 自動帶入",
    bullets=[
        "# Picker 三層",
        "1. 專區（橫流溪 / 烏石坑）",
        "2. 承租人（依專區自動篩選）",
        "3. 作業單元（依承租人列出 + 契約面積/樹種摘要）",
        "# 完成選擇後自動發生",
        "申請人姓名 → 自動覆蓋為承租人姓名（v2.11.63）",
        "林班 / 地號、面積、株數 → 自動帶入",
        "迷你地圖 → 紅色高亮目標 + 自動 zoom",
        "下方 7 欄系統識別摘要",
        "# 仍可手動覆蓋任何欄位",
    ],
    image="S06-Picker三層連動完成.png",
    image_caption="picker 完成、mini-map 紅色高亮、申請人姓名已帶入",
    image_width=4.5, image_left=8.2, image_top=1.4)

# =====================================================
# Slide 16: 表單下半部
# =====================================================
add_content_slide(prs, "表單下半部 — 自動計算與必填欄位",
    bullets=[
        "# 自動連動計算",
        ("修枝株數", "面積 × 1800 株/ha（合作社契約假設密度）"),
        "面積變更時即時更新株數",
        "可手動覆蓋",
        "# 必填欄位（紅色 *）",
        "聯絡方式（電話／email）",
        "修枝起日",
        "修枝迄日（須晚於起日）",
        "# 送出選項",
        ("💾 儲存草稿", "可日後繼續編輯，不配發文號"),
        ("📤 送出申請", "即時配發發文字號「（修）字第 NNN 號」"),
    ],
    image="S07-申請表單下半部.png",
    image_caption="表單下半部欄位",
    image_width=4.5, image_left=8.2, image_top=1.4)

# =====================================================
# Slide 17: 申請公文稿 — 第 1 頁
# =====================================================
add_big_image_slide(prs, "申請公文稿（第 1 頁）— 公文體 4 點 + 表列",
    image="S08-申請公文稿-第1頁.png",
    image_caption="正式公文「函」格式：受文者 / 發文字號 / 公文體說明 4 點 / 表列 8 欄申請明細",
    image_width=8.5)

# =====================================================
# Slide 18: 申請公文稿 — 第 2 頁 SVG
# =====================================================
add_big_image_slide(prs, "申請公文稿（第 2 頁）— 作業位置 SVG 示意圖",
    image="S09-申請公文稿-第2頁SVG.png",
    image_caption="自製 inline SVG：紅目標 + 灰 context + N 北方箭頭 + meta 行 + 申請人簽章 + 分署收件欄",
    image_width=8.5)

# =====================================================
# 第 4 段 分隔頁
# =====================================================
add_section_divider(prs, "4", "分署承辦審核與核准", "15 - 20 分鐘")

# =====================================================
# Slide 20: 審核清單
# =====================================================
add_content_slide(prs, "修枝審核分頁（分署承辦視角）",
    bullets=[
        "# 權限聲明（副標）",
        "「審核者僅能核准／要求補件／駁回，不可變更申請內容（權限由 Firestore 規則鎖定）」",
        "# 分區",
        "待審核（NN）：尚未審核之案件",
        "已處理（NN）：已核准/駁回之歷史案件",
        "# 三按鈕",
        "✅ 核准（綠）",
        "📋 要求補件（橙）",
        "✕ 駁回（紅）",
    ],
    image="S10-修枝審核清單.png",
    image_caption="修枝審核清單（含三筆已處理案件）",
    image_width=6.0, image_left=6.8, image_top=1.4)

# =====================================================
# Slide 21: 核准 Modal (v2.11.69)
# =====================================================
add_content_slide(prs, "核准 Modal（v2.11.69 拍板：不預估採收量）",
    bullets=[
        "# Modal 內容",
        "申請人資訊摘要（read-only）",
        ("灰底宣告框", "「本案核准後即許可進行修枝作業；所採取葉片重量由林農於 ForestMRV 線上系統實際回報，不設預估上限」"),
        "效期起日 / 迄日",
        "審核附註（選填）",
        "# 為什麼不設「核准採收量」上限？",
        "土肉桂每塊地年齡不同、每株產出葉量差異大",
        "事先預估上限反而誤導承辦與林農",
        "改採「事後實際回報」設計",
        "合作社共同銷售依實際量結算",
    ],
    image="S11-核准Modal-v2.11.69.png",
    image_caption="新版核准 Modal（v2.11.69）",
    image_width=4.5, image_left=8.2, image_top=1.4)

# =====================================================
# Slide 22: 法定許可文號 Transaction
# =====================================================
_s = new_diagram_slide("法定許可文號自動配發")
_ny, _nh, _ay, _ah = 2.05, 1.7, 2.55, 0.7
add_node(_s, 0.55, _ny, 2.6, _nh, [
    ("Firestore", 16, False, STONE_100),
    ("runTransaction", 18, True, WHITE),
    ("原子遞增", 16, False, STONE_100)], fill=FOREST_700)
add_arrow(_s, 3.15, _ay, 0.6, _ah)
add_node(_s, 3.75, _ny, 2.6, _nh, [
    ("連號文號配發", 18, True, WHITE),
    ("林保中-土肉桂修枝", 13, False, STONE_100),
    ("-115-001", 15, True, AMBER_500)], fill=FOREST_700)
add_arrow(_s, 6.35, _ay, 0.6, _ah)
add_node(_s, 6.95, _ny, 2.6, _nh, [
    ("申請狀態推進", 18, True, WHITE),
    ("submitted", 14, False, STONE_100),
    ("→ approved", 15, True, AMBER_500)], fill=FOREST_700)
add_arrow(_s, 9.55, _ay, 0.6, _ah)
add_node(_s, 10.15, _ny, 2.6, _nh, [
    ("林農端即時可見", 18, True, WHITE),
    ("卡片轉綠『已核准』", 13, False, STONE_100),
    ("出現採取許可單", 13, False, STONE_100)], fill=FOREST_800)
add_notes_box(_s, 0.6, 4.35, 12.1, 2.4, [
    "# 為什麼用 Transaction？",
    "確保連號、不可重號、不可並寫 — 多承辦同時核准也安全",
    "文號格式：林保中-土肉桂修枝-{民國年}-{流水號}（例 115-001）",
    "核准同步寫入效期與審核附註，並以 Toast 通知承辦",
], size=20)

# =====================================================
# Slide 23: 採取許可單
# =====================================================
add_big_image_slide(prs, "採取許可單 — 列印版（含 v2.11.69 宣告）",
    image="S12-採取許可單-列印版.png",
    image_caption="採取許可單公文 box + 一句宣告（dashed border-top）「本許可只許可修枝作業；所採取葉片重量以實際回報為準、不設預估上限」",
    image_width=8.5)

# =====================================================
# 第 5 段 分隔頁
# =====================================================
add_section_divider(prs, "5", "林農回報與結案", "20 - 25 分鐘")

# =====================================================
# Slide 25: 回報分頁
# =====================================================
add_content_slide(prs, "葉片採收回報及結案分頁",
    bullets=[
        "# 業務說明",
        "核准效期內進行修枝作業並採收葉片",
        "回報為法定義務，未回報不可結案",
        "# 回報設計（v2.11.69 簡化）",
        ("無對照、無達成率", "純累計顯示，不比對核准量"),
        ("可分批多次", "如 5/30 採 10 kg、6/15 採 8 kg、累計 18 kg"),
        "# 回報欄位",
        "採收日期、實際鮮葉重(kg) — 必填",
        "乾燥後重、含水率、批次標示、備註 — 選填",
        "# 自動推進狀態",
        "首筆回報時 approved → harvesting（修枝作業中）",
    ],
    image="S14-葉片採收回報及結案.png",
    image_caption="採收回報及結案分頁",
    image_width=6.0, image_left=6.8, image_top=1.4)

# =====================================================
# Slide 26: 結案閘門
# =====================================================
_s = new_diagram_slide("回報完畢並結案 — 結案閘門（v2.11.39）")
# 起點
add_node(_s, 0.55, 2.55, 2.7, 1.5, [
    ("累計回報量", 18, True, WHITE),
    ("totalLogged_kg", 14, False, STONE_100)], fill=STONE_700)
add_arrow(_s, 3.3, 3.05, 0.55, 0.55)
# 決策菱形
add_node(_s, 3.95, 2.2, 3.0, 2.2, [
    ("回報量 > 0 ？", 20, True, WHITE)], fill=BLUE_700,
    shape=MSO_SHAPE.DIAMOND)
# 是 → 可結案
add_arrow(_s, 7.05, 2.35, 0.6, 0.55)
add_node(_s, 7.7, 1.7, 4.9, 1.35, [
    ("✓ 是 → 可結案", 18, True, WHITE),
    ("狀態 harvesting → completed（已結案）", 14, False, STONE_100)], fill=FOREST_700)
# 否 → 擋下
add_arrow(_s, 7.05, 3.7, 0.6, 0.55)
add_node(_s, 7.7, 3.55, 4.9, 1.35, [
    ("✗ 否 → 擋下", 18, True, WHITE),
    ("提示先回報採收量，不可結案", 14, False, STONE_100)], fill=RED_600)
add_notes_box(_s, 0.6, 5.15, 12.1, 1.9, [
    "# 雙重把關",
    "法定義務：未回報（totalLogged_kg = 0）不可結案 — Firestore Rules + 客戶端同步擋",
    "結案後不可再回報、不可再編輯；資料固定供分署查核與合作社彙整",
], size=19)

# =====================================================
# 第 6 段 分隔頁
# =====================================================
add_section_divider(prs, "6", "合作社彙整與 Q&A", "25 - 30 分鐘")

# =====================================================
# Slide 28: 合作社彙整
# =====================================================
add_big_image_slide(prs, "葉片採收彙整（合作社視角）",
    image="S13-葉片採收彙整.png",
    image_caption="三區塊：📋 專區總覽 + 🤝 依林農三欄表（無達成率欄）+ 📋 各申請案唯讀",
    image_width=10.5)

# =====================================================
# Slide 29: 共同銷售規劃
# =====================================================
_s = new_diagram_slide("共同銷售規劃（合作社彙整三區塊）")
add_node(_s, 0.55, 1.5, 5.6, 1.45, [
    ("區塊一：專區總覽", 18, True, WHITE),
    ("申請案總數 · 狀態分布 · 已回報採收量總計", 13, False, STONE_100)], fill=FOREST_700)
add_node(_s, 0.55, 3.1, 5.6, 1.45, [
    ("區塊二：依林農彙整", 18, True, WHITE),
    ("三欄 林農／案件數／已回報量(kg)，降序，無達成率", 13, False, STONE_100)], fill=FOREST_700)
add_node(_s, 0.55, 4.7, 5.6, 1.45, [
    ("區塊三：各申請案（唯讀）", 18, True, WHITE),
    ("可預覽每筆採取許可單", 13, False, STONE_100)], fill=FOREST_700)
add_arrow(_s, 6.35, 3.35, 0.95, 0.9)
add_node(_s, 7.5, 2.3, 5.2, 2.6, [
    ("📥 匯出 Excel", 24, True, WHITE),
    ("", 8, False, WHITE),
    ("供共同銷售前資料整理", 16, False, STONE_100),
    ("通路議價 · 結算分潤", 16, False, STONE_100)], fill=FOREST_800)
add_notes_box(_s, 7.5, 5.15, 5.2, 1.4, [
    "v2.11.41 拍板：純累計顯示、不設達成率",
], size=17)

# =====================================================
# Slide 30: 五個角色重點回顧
# =====================================================
_s = new_diagram_slide("結語：五個角色的協作鏈")
_cy, _ch, _cay, _cah = 1.65, 2.2, 2.4, 0.7
add_node(_s, 0.35, _cy, 2.25, _ch, [
    ("1. admin", 16, True, AMBER_500),
    ("系統管理員", 16, True, WHITE),
    ("開案 · 邊界", 12, False, STONE_100),
    ("邀請成員", 12, False, STONE_100)], fill=FOREST_700)
add_arrow(_s, 2.6, _cay, 0.35, _cah)
add_node(_s, 2.95, _cy, 2.25, _ch, [
    ("2. surveyor", 16, True, AMBER_500),
    ("林農", 16, True, WHITE),
    ("修枝申請", 12, False, STONE_100),
    ("+ 紙本公文", 12, False, STONE_100)], fill=FOREST_700)
add_arrow(_s, 5.2, _cay, 0.35, _cah)
add_node(_s, 5.55, _cy, 2.25, _ch, [
    ("3. harvest_", 15, True, AMBER_500),
    ("authority｜分署", 15, True, WHITE),
    ("審核 · 核准", 12, False, STONE_100),
    ("配發法定文號", 12, False, STONE_100)], fill=BLUE_700)
add_arrow(_s, 7.8, _cay, 0.35, _cah)
add_node(_s, 8.15, _cy, 2.25, _ch, [
    ("4. surveyor", 16, True, AMBER_500),
    ("林農", 16, True, WHITE),
    ("修枝採葉", 12, False, STONE_100),
    ("回報 · 結案", 12, False, STONE_100)], fill=FOREST_700)
add_arrow(_s, 10.4, _cay, 0.35, _cah)
add_node(_s, 10.75, _cy, 2.25, _ch, [
    ("5. coop", 16, True, AMBER_500),
    ("合作社", 16, True, WHITE),
    ("唯讀彙整", 12, False, STONE_100),
    ("匯出 Excel", 12, False, STONE_100)], fill=STONE_700)
add_notes_box(_s, 0.6, 4.55, 12.1, 2.3, [
    "# v2.11.69 業務拍板回顧",
    "許可只許可修枝作業，不預估／不限制採收量上限",
    "採收量以實際回報為準、不對照、無達成率",
    "「事後實際彙整」設計符合土肉桂業務本質",
], size=20)

# =====================================================
# Slide 31: Q&A 結束
# =====================================================
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)
add_full_bleed_bg(slide, IMG_BG_SECTION, FOREST_900, 62, FOREST_900)

tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.8))
tf = tb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Q & A"
set_text(run, size=100, bold=True, color=AMBER_500)

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)
run2 = p2.add_run()
run2.text = "歡迎提問"
set_text(run2, size=40, color=WHITE)

tb_info = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.8))
tf_info = tb_info.text_frame
p_info = tf_info.paragraphs[0]
p_info.alignment = PP_ALIGN.CENTER
run_info = p_info.add_run()
run_info.text = "土肉桂專區葉片收穫監測會計系統 │ ForestMRV v2.11.69"
set_text(run_info, size=18, color=STONE_300)

# 存檔
prs.save(OUT)
print(f"OK: {OUT}")
print(f"Size: {OUT.stat().st_size / 1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
