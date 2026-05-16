# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Emu

f = r'05-土肉桂採收許可-520說明會簡報.pptx'
pr = Presentation(f)
SW = pr.slide_width; SH = pr.slide_height
print('slides =', len(pr.slides), ' size =', round(Emu(SW).inches,2), 'x', round(Emu(SH).inches,2))
alltext = []
for i, sl in enumerate(pr.slides, 1):
    texts = []
    oob = []
    for sh in sl.shapes:
        t = ''
        if sh.has_text_frame:
            t = '\n'.join(p.text for p in sh.text_frame.paragraphs).strip()
        elif sh.has_table:
            rows = []
            for r in sh.table.rows:
                rows.append(' | '.join(c.text for c in r.cells))
            t = '[TABLE] ' + ' || '.join(rows)
        if t:
            texts.append(t)
        try:
            l = Emu(sh.left).inches; tp = Emu(sh.top).inches
            w = Emu(sh.width).inches; h = Emu(sh.height).inches
            if l < -0.02 or tp < -0.02 or l + w > round(Emu(SW).inches,2) + 0.05 or tp + h > round(Emu(SH).inches,2) + 0.05:
                oob.append(f'  !! OOB shape L={l:.2f} T={tp:.2f} R={l+w:.2f} B={tp+h:.2f} "{(t[:24] if t else sh.shape_type)}"')
        except Exception:
            pass
    body = '\n'.join(texts)
    alltext.append(body)
    note = ''
    if sl.has_notes_slide:
        note = sl.notes_slide.notes_text_frame.text.strip()
    print('\n===== Slide', i, '=====')
    print(body)
    print('--- notes:', (note[:70] + ('...' if len(note) > 70 else '')) if note else '(無)')
    for o in oob:
        print(o)

joined = '\n'.join(alltext)
print('\n===== 4 必含項檢查 =====')
for k in ['申請公文', '林產物採取申請函', '受文者', '申請表單', '欄位', '許可文號', '林保中-土肉桂採葉', '合作社', '已回報', '達成率', '168']:
    print(('  OK ' if k in joined else '  MISSING ') + k)
print('\n===== placeholder 掃描 =====')
import re
bad = re.findall(r'(?i)lorem|ipsum|\bTODO\b|\[insert|xxxx', joined)
print('  ', bad if bad else '無 placeholder 殘留')
