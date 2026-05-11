"""Inspect placeholders for H01, M08, P04 in workshop materials."""
import sys
from pathlib import Path
from docx import Document
from pptx import Presentation

BASE = Path(r"H:\我的雲端硬碟\2026 forestry_RS\reports\2026-05-mid-demo")
TARGETS = ["H01", "M08", "P04"]


def scan_docx(path: Path, name: str):
    print(f"\n=== {name} :: {path.name} ===")
    doc = Document(str(path))
    for pi, p in enumerate(doc.paragraphs):
        t = p.text
        for tag in TARGETS:
            if tag in t:
                print(f"  [para#{pi}] {t!r}")
    for ti, table in enumerate(doc.tables):
        for ri, row in enumerate(table.rows):
            for ci, cell in enumerate(row.cells):
                t = cell.text
                for tag in TARGETS:
                    if tag in t:
                        print(f"  [tbl#{ti} r{ri} c{ci}] {t!r}")


def scan_pptx(path: Path, name: str):
    print(f"\n=== {name} :: {path.name} ===")
    prs = Presentation(str(path))
    for si, slide in enumerate(prs.slides, start=1):
        for shi, sh in enumerate(slide.shapes):
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            for tag in TARGETS:
                if tag in t:
                    print(f"  [slide#{si} shape#{shi} name={sh.name!r}] {t!r}")


scan_pptx(BASE / "01-簡報-ForestMRV系統介紹.pptx", "PPT")
scan_docx(BASE / "02-操作手冊-ForestMRV.docx", "MANUAL")
scan_docx(BASE / "03-現場練習指南-下午hands-on.docx", "HANDS-ON")
