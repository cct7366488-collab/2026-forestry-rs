"""Embed 3 screenshots (H01, M08, P04) into workshop materials.

Strategy:
- PPT slide 19 shape 'Text 6': remove text, insert P04 picture at same bounding box.
- MANUAL table#25 r0 c0: clear cell text, add picture run inside the cell.
- HANDS-ON para#22: clear paragraph text, add picture run.

Each file is backed up to *.bak-YYYYMMDD-HHMMSS before write.
"""
from __future__ import annotations

import shutil
import sys
from datetime import datetime
from pathlib import Path

from docx import Document
from docx.shared import Cm, Emu
from PIL import Image
from pptx import Presentation
from pptx.util import Emu as PptxEmu

BASE = Path(r"H:\我的雲端硬碟\2026 forestry_RS\reports\2026-05-mid-demo")
IMG = BASE / "_build" / "images"

P04 = IMG / "P04-新專案form.png"
M08 = IMG / "M08-Chrome安裝圖示.png"
H01 = IMG / "H01-章節導引示意.png"

assert P04.exists() and M08.exists() and H01.exists(), "missing image(s)"

PPT_PATH = BASE / "01-簡報-ForestMRV系統介紹.pptx"
MAN_PATH = BASE / "02-操作手冊-ForestMRV.docx"
HOL_PATH = BASE / "03-現場練習指南-下午hands-on.docx"

STAMP = datetime.now().strftime("%Y%m%d-%H%M%S")


def backup(p: Path) -> Path:
    bak = p.with_suffix(p.suffix + f".bak-{STAMP}")
    shutil.copy2(p, bak)
    return bak


def image_size_emu(img_path: Path, max_w_emu: int, max_h_emu: int) -> tuple[int, int]:
    """Return (width, height) in EMU, fit-inside max box, preserve aspect."""
    with Image.open(img_path) as im:
        w_px, h_px = im.size
    ratio = w_px / h_px
    box_ratio = max_w_emu / max_h_emu
    if ratio > box_ratio:
        w = max_w_emu
        h = int(w / ratio)
    else:
        h = max_h_emu
        w = int(h * ratio)
    return w, h


# ----- PPT: slide 19 'Text 6' -> P04 -----
def embed_ppt():
    print(f"[PPT ] backup -> {backup(PPT_PATH).name}")
    prs = Presentation(str(PPT_PATH))
    slide = prs.slides[18]  # slide 19 (1-indexed)
    target = None
    for sh in slide.shapes:
        if sh.name == "Text 6":
            target = sh
            break
    if target is None:
        raise SystemExit("PPT: Text 6 shape not found")
    left, top, width, height = target.left, target.top, target.width, target.height
    # remove the text shape
    sp = target._element
    sp.getparent().remove(sp)
    # add picture fit inside the original bbox
    w, h = image_size_emu(P04, width, height)
    # center inside original bbox
    new_left = left + (width - w) // 2
    new_top = top + (height - h) // 2
    pic = slide.shapes.add_picture(str(P04), new_left, new_top, width=w, height=h)
    print(f"[PPT ] slide 19 picture inserted, size={pic.width / 914400:.2f}x{pic.height / 914400:.2f} in")
    prs.save(str(PPT_PATH))


# ----- DOCX helpers -----
def replace_paragraph_with_picture(paragraph, img_path: Path, width_cm: float):
    """Clear paragraph runs and insert a single picture run."""
    # remove all existing runs in the paragraph
    for r in list(paragraph.runs):
        r._element.getparent().remove(r._element)
    # also strip remaining text in pPr child? python-docx runs cover it.
    run = paragraph.add_run()
    run.add_picture(str(img_path), width=Cm(width_cm))


def embed_manual():
    print(f"[MAN ] backup -> {backup(MAN_PATH).name}")
    doc = Document(str(MAN_PATH))
    table = doc.tables[25]
    cell = table.rows[0].cells[0]
    # find paragraph containing the placeholder
    target_para = None
    for p in cell.paragraphs:
        if "M08" in p.text and "screenshot" in p.text:
            target_para = p
            break
    if target_para is None:
        raise SystemExit("MAN: M08 placeholder paragraph not found")
    # remove the OTHER paragraphs in cell that are duplicate caption ('Chrome 「安裝」圖示位置')
    keep_caption = None
    for p in cell.paragraphs:
        if p is target_para:
            continue
        if p.text.strip() == "Chrome 「安裝」圖示位置":
            keep_caption = p  # keep this as caption beneath
            break
    replace_paragraph_with_picture(target_para, M08, width_cm=14.0)
    print("[MAN ] table#25 cell M08 picture inserted")
    doc.save(str(MAN_PATH))


def embed_handson():
    print(f"[HOL ] backup -> {backup(HOL_PATH).name}")
    doc = Document(str(HOL_PATH))
    target_para = None
    for p in doc.paragraphs:
        if "H01" in p.text and "章節導引" in p.text:
            target_para = p
            break
    if target_para is None:
        raise SystemExit("HOL: H01 placeholder paragraph not found")
    replace_paragraph_with_picture(target_para, H01, width_cm=14.0)
    print("[HOL ] para H01 picture inserted")
    doc.save(str(HOL_PATH))


embed_ppt()
embed_manual()
embed_handson()
print("\nDONE.")
